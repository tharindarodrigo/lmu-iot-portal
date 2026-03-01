#!/usr/bin/env python3
"""Generate a clean, presentation-ready PPTX deck from project content."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
IMAGE_DIR = DOCS_DIR / "report-images"
OUTPUT_PPTX = DOCS_DIR / "presentation-demo-14min-clean.pptx"
OUTPUT_PDF = DOCS_DIR / "presentation-demo-14min-clean.pdf"

SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

COLORS: dict[str, RGBColor] = {
    "background": RGBColor(248, 250, 252),
    "surface": RGBColor(255, 255, 255),
    "line": RGBColor(203, 213, 225),
    "title": RGBColor(15, 23, 42),
    "body": RGBColor(30, 41, 59),
    "muted": RGBColor(71, 85, 105),
    "accent": RGBColor(20, 184, 166),
}

IMAGES: dict[str, Path] = {
    "architecture": IMAGE_DIR / "figure2-end-to-end-architecture.png",
    "automation": IMAGE_DIR / "figure3-automation-dag.png",
    "dashboard": IMAGE_DIR / "figure4-iot-dashboard.png",
    "control": IMAGE_DIR / "figure5-device-control-feedback.png",
    "deployment": IMAGE_DIR / "deployment-architecture-aks.png",
}


def set_background(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_WIDTH),
        Inches(SLIDE_HEIGHT),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["background"]
    shape.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.65))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS["title"]

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.84), Inches(12.1), Inches(0.35))
        sub_tf = subtitle_box.text_frame
        sub_tf.clear()
        sp = sub_tf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = Pt(15)
        sp.font.color.rgb = COLORS["muted"]

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(1.22),
        Inches(12.1),
        Inches(0.03),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS["accent"]
    divider.line.fill.background()


def add_speaker_notes(slide, notes: str) -> None:
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = notes


def add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str], size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = COLORS["body"]
        paragraph.space_after = Pt(10)


def add_text(slide, left: float, top: float, width: float, height: float, text: str, size: int = 18, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS["body"]


def add_image_contain(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["surface"]
    frame.line.color.rgb = COLORS["line"]

    if not image_path.exists():
        return

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    target_ratio = width / height
    image_ratio = image_width / image_height

    if image_ratio > target_ratio:
        draw_width = width - 0.35
        draw_height = draw_width / image_ratio
    else:
        draw_height = height - 0.35
        draw_width = draw_height * image_ratio

    draw_left = left + (width - draw_width) / 2
    draw_top = top + (height - draw_height) / 2

    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_left),
        Inches(draw_top),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )


def add_two_column_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets: list[str],
    image_key: str,
    image_caption: str,
    notes: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, subtitle)
    add_bullets(slide, 0.75, 1.55, 5.6, 5.55, bullets, size=18)
    add_image_contain(slide, IMAGES[image_key], 6.45, 1.55, 6.1, 4.95)
    add_text(slide, 6.52, 6.6, 5.95, 0.5, image_caption, size=13)
    add_speaker_notes(slide, notes)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Cloud IoT Energy Monitoring and Automation Platform", "Presentation and live demonstration")
    add_text(
        slide,
        0.8,
        2.0,
        11.9,
        1.0,
        "End-to-end system: telemetry ingestion, 15-minute rule evaluation, automated control, and auditable evidence.",
        size=24,
        bold=True,
    )
    add_bullets(
        slide,
        0.8,
        3.25,
        11.4,
        2.8,
        [
            "Single integrated platform boundary",
            "Simulated energy device for live proof of behavior",
            "Focus on operational correctness and traceability",
        ],
        size=20,
    )
    add_speaker_notes(
        slide,
        (
            "Today we are presenting our cloud IoT energy monitoring and automation platform.\n"
            "The key point is that this system does not stop at visualization. It takes telemetry, evaluates a rule, sends a control command, and verifies the response.\n"
            "For the demo, we will use a simulated energy device so you can see the same end-to-end flow used for real devices."
        ),
    )

    # Slide 2: Problem and objective
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Problem and Objective", "From passive monitoring to automated action")
    add_bullets(
        slide,
        0.75,
        1.65,
        12.0,
        5.6,
        [
            "Many IoT deployments stop at dashboards and manual intervention",
            "Operational reaction is delayed when humans must inspect every change",
            "This project implements a deterministic telemetry-to-action loop",
            "Goal: trigger control behavior based on validated 15-minute energy conditions",
            "All outcomes must be traceable through run and command records",
        ],
        size=20,
    )
    add_speaker_notes(
        slide,
        (
            "The problem we solved is operational delay.\n"
            "In many IoT setups, teams watch dashboards and manually react. Our objective was to make this deterministic.\n"
            "We implemented a 15-minute energy rule that triggers control only when the condition is met, and we record every decision and outcome for auditability."
        ),
    )

    # Slide 3: Architecture
    add_two_column_slide(
        prs,
        title="End-to-End Architecture",
        subtitle="A single telemetry substrate for ingestion, automation, control, and reporting",
        bullets=[
            "Devices publish telemetry over MQTT-style topic contracts",
            "Ingestion validates and persists canonical telemetry records",
            "Automation evaluates conditions and dispatches control commands",
            "Dashboard and reporting consume the same consistent data source",
        ],
        image_key="architecture",
        image_caption="Architecture overview",
        notes=(
            "This diagram shows the full path.\n"
            "Devices publish telemetry through MQTT-style topics, ingestion validates and stores canonical records, and automation evaluates rules before dispatching control commands.\n"
            "Dashboards and reports read from the same data source, so monitoring, control, and evidence stay consistent."
        ),
    )

    # Slide 4: Firmware approach
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Firmware Strategy", "Reusable ESP32 runtime with modular device logic")
    add_bullets(
        slide,
        0.75,
        1.65,
        5.6,
        5.4,
        [
            "Common runtime handles WiFi, MQTT, presence, and reconnect behavior",
            "RGB actuator module implements control + state acknowledgment",
            "PZEM meter module implements RS485 Modbus telemetry polling",
            "Read-only meter path in v1 avoids unsafe write-side effects",
        ],
        size=18,
    )
    add_image_contain(slide, IMAGES["dashboard"], 6.45, 1.55, 6.1, 4.95)
    add_text(slide, 6.52, 6.6, 5.95, 0.5, "Telemetry dashboard view", size=13)
    add_speaker_notes(
        slide,
        (
            "On the firmware side, we use a shared runtime and separate device modules.\n"
            "The shared runtime handles WiFi, MQTT, presence, and reconnect behavior. The RGB module handles control and acknowledgments, and the PZEM module handles Modbus telemetry polling.\n"
            "This approach makes onboarding new devices faster and keeps v1 meter behavior safely read-only."
        ),
    )

    # Slide 5: Automation logic
    add_two_column_slide(
        prs,
        title="Automation Logic",
        subtitle="15-minute window rule with deterministic branching",
        bullets=[
            "Rule: MAX(total_energy_kwh) - MIN(total_energy_kwh) over 15 minutes",
            "Threshold pass branch dispatches an RGB alert command",
            "Threshold fail branch correctly avoids control side effects",
            "Command lifecycle is reconciled with device feedback for auditability",
        ],
        image_key="automation",
        image_caption="Automation DAG and condition flow",
        notes=(
            "This is the exact rule in our demo: the 15-minute consumption delta is MAX(total_energy_kwh) minus MIN(total_energy_kwh).\n"
            "If that value crosses the threshold, the workflow dispatches an RGB alert command. If it does not cross, no command is sent.\n"
            "Both branches are explicit and every step is logged through run records and command lifecycle states."
        ),
    )

    # Slide 6: Demo walkthrough (no timestamps)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Demo Walkthrough", "Simulation-first proof of end-to-end behavior")
    add_bullets(
        slide,
        0.75,
        1.65,
        5.6,
        5.4,
        [
            "Start from normal baseline with no active alert",
            "Publish simulated meter telemetry with rising energy values",
            "Observe real-time dashboard updates and rule evaluation",
            "Verify command dispatch and device acknowledgment",
            "Show reporting evidence and a safe negative-path case",
        ],
        size=18,
    )
    add_image_contain(slide, IMAGES["control"], 6.45, 1.55, 6.1, 4.95)
    add_text(slide, 6.52, 6.6, 5.95, 0.5, "Control lifecycle and feedback reconciliation", size=13)
    add_speaker_notes(
        slide,
        (
            "Now for the demo flow.\n"
            "We start from a normal baseline, publish simulated meter telemetry with rising energy values, and observe real-time updates and rule evaluation.\n"
            "Then we verify command dispatch and acknowledgment, and finish with reporting evidence plus a safe negative-path case where no command is triggered."
        ),
    )

    # Slide 7: Deployment and scaling
    add_two_column_slide(
        prs,
        title="Deployment and Scalability",
        subtitle="Cloud-agnostic architecture with Kubernetes-ready separation",
        bullets=[
            "Workloads can scale independently across ingestion, automation, and reporting",
            "Containerized service boundaries support managed Kubernetes platforms",
            "Queue-backed processing smooths spikes in telemetry and export requests",
            "Contract-driven onboarding supports future device expansion",
        ],
        image_key="deployment",
        image_caption="Portable deployment model",
        notes=(
            "This architecture is designed to scale beyond the classroom demo.\n"
            "Ingestion, automation, and reporting workloads can scale independently, and the containerized design stays cloud-agnostic and Kubernetes-ready.\n"
            "That gives us a practical path from prototype behavior to production-style operations."
        ),
    )

    # Slide 8: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Conclusion", "Full-loop IoT operations demonstrated")
    add_bullets(
        slide,
        0.85,
        1.8,
        11.6,
        4.8,
        [
            "Telemetry is converted into deterministic automated control",
            "The platform provides clear operational and audit evidence",
            "Architecture and firmware patterns are reusable for future growth",
            "Result: ingest, decide, act, and verify within one coherent system",
        ],
        size=22,
    )
    add_text(slide, 0.85, 6.55, 11.8, 0.5, "Q&A", size=30, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(
        slide,
        (
            "To conclude, we demonstrated a full-loop IoT workflow: ingest, decide, act, and verify.\n"
            "The platform combines automation with traceability, so outcomes are operationally useful and auditable.\n"
            "The same architecture and firmware pattern can be reused for future devices and larger deployments. Thank you, and we welcome your questions."
        ),
    )

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    return OUTPUT_PPTX


def build_pdf(pptx_path: Path) -> Path:
    import aspose.slides as slides

    with slides.Presentation(str(pptx_path)) as presentation:
        presentation.save(str(OUTPUT_PDF), slides.export.SaveFormat.PDF)

    return OUTPUT_PDF


if __name__ == "__main__":
    pptx = build_pptx()
    print(f"Generated PPTX: {pptx}")

    try:
        pdf = build_pdf(pptx)
        print(f"Generated PDF: {pdf}")
    except Exception as exc:  # pragma: no cover
        print(f"PDF generation failed: {exc}")
